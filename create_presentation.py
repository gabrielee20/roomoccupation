#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script per creare una presentazione PowerPoint dettagliata
sul progetto di predizione dell'occupazione della stanza
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs, title, subtitle=""):
    """Aggiunge una slide con titolo"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle
    
    return slide

def add_content_slide(prs, title, content_points):
    """Aggiunge una slide con titolo e contenuto bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        # Gestisce indentazione (se inizia con "  -")
        if point.startswith("  "):
            p.level = 1
            point = point.strip()
        
        p.text = point.lstrip("- ")
        p.font.size = Pt(18)
    
    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Aggiunge una slide con immagine"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Aggiungi titolo
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Aggiungi immagine se esiste
    if os.path.exists(image_path):
        left = Inches(1)
        top = Inches(1.5)
        
        # Calcola dimensioni mantenendo aspect ratio
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(8))
    
    # Aggiungi caption se presente
    if caption:
        left = Inches(1)
        top = Inches(6.5)
        width = Inches(8)
        height = Inches(0.5)
        
        caption_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = caption_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Crea la presentazione completa"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Titolo principale
    add_title_slide(
        prs,
        "Predizione dell'Occupazione della Stanza",
        "Progetto di Machine Learning con Random Forest\n\nAnalisi, Implementazione e Risultati"
    )
    
    # Slide 2: Indice
    add_content_slide(
        prs,
        "Indice della Presentazione",
        [
            "Introduzione e Obiettivi",
            "Dataset e Caratteristiche",
            "Analisi Esplorativa dei Dati",
            "Visualizzazione delle Features",
            "Metodologia di Classificazione",
            "Ottimizzazione degli Iper-parametri",
            "Training del Modello",
            "Valutazione delle Performance",
            "Feature Importance",
            "Conclusioni e Risultati Finali"
        ]
    )
    
    # Slide 3: Introduzione
    add_content_slide(
        prs,
        "Introduzione al Progetto",
        [
            "Obiettivo: Predire l'occupazione di una stanza basandosi su sensori ambientali",
            "Approccio: Machine Learning con algoritmo Random Forest",
            "Dataset: Dati da sensori di temperatura, umidità, luce, CO2 e umidità",
            "Applicazioni pratiche:",
            "  - Gestione energetica intelligente degli edifici",
            "  - Ottimizzazione del condizionamento",
            "  - Sicurezza e monitoraggio degli spazi"
        ]
    )
    
    # Slide 4: Obiettivi Specifici
    add_content_slide(
        prs,
        "Obiettivi Specifici del Progetto",
        [
            "Analizzare e comprendere il dataset di occupazione",
            "Identificare le features più rilevanti per la predizione",
            "Implementare un modello di classificazione Random Forest",
            "Ottimizzare gli iper-parametri tramite Grid Search",
            "Valutare le performance con multiple metriche",
            "Validare il modello su dataset separati"
        ]
    )
    
    # Slide 5: Dataset
    add_content_slide(
        prs,
        "Dataset Utilizzati",
        [
            "Tre file di dataset:",
            "  - datatraining.txt: dati per il training iniziale",
            "  - datatest.txt: primo set di test",
            "  - datatest2.txt: validazione finale separata",
            "Features disponibili:",
            "  - Temperature (Temperatura)",
            "  - Humidity (Umidità)",
            "  - Light (Luminosità)",
            "  - CO2 (Anidride Carbonica)",
            "  - HumidityRatio (Rapporto di Umidità)",
            "Target: Occupancy (0 = non occupata, 1 = occupata)"
        ]
    )
    
    # Slide 6: Struttura dei Dati
    add_content_slide(
        prs,
        "Caratteristiche del Dataset",
        [
            "Dati temporali con timestamp",
            "Campionamento continuo nel tempo",
            "Divisione strategica dei dati:",
            "  - 75% per training",
            "  - 25% per testing (subsampling stratificato)",
            "Mantenimento del dataset test2 per validazione finale",
            "Separazione features (X) e label (y) per entrambi i set"
        ]
    )
    
    # Slide 7: Analisi Esplorativa - Giorni
    if os.path.exists("plots/campioni_con_presenza_per_giorno.png"):
        add_image_slide(
            prs,
            "Distribuzione dell'Occupazione per Giorno",
            "plots/campioni_con_presenza_per_giorno.png",
            "Analisi della presenza per giorno della settimana"
        )
    else:
        add_content_slide(
            prs,
            "Distribuzione dell'Occupazione per Giorno",
            [
                "Analisi della presenza per giorno della settimana",
                "Osservazioni principali:",
                "  - Alta presenza: Lunedì, Giovedì, Venerdì",
                "  - Nessuna presenza: Sabato e Domenica",
                "  - Pattern tipico di ambiente lavorativo/scolastico"
            ]
        )
    
    # Slide 8: Heatmap Oraria
    if os.path.exists("plots/heatmap_campioni_per_ora.png"):
        add_image_slide(
            prs,
            "Heatmap: Campioni per Ora del Giorno",
            "plots/heatmap_campioni_per_ora.png",
            "Distribuzione oraria dell'occupazione"
        )
    else:
        add_content_slide(
            prs,
            "Analisi Temporale Oraria",
            [
                "Distribuzione dell'occupazione nelle diverse ore",
                "Pattern orari identificati",
                "Correlazione con orari lavorativi/scolastici"
            ]
        )
    
    # Slide 9: Andamento Temporale Features
    if os.path.exists("plots/andamento_temporale_features.png"):
        add_image_slide(
            prs,
            "Andamento Temporale delle Features",
            "plots/andamento_temporale_features.png",
            "Visualizzazione dell'andamento di ogni feature nel tempo (rosso = presenza)"
        )
    else:
        add_content_slide(
            prs,
            "Andamento Temporale delle Features",
            [
                "Visualizzazione usando small multiples",
                "Evidenziazione in rosso quando qualcuno è presente",
                "Separazione dei giorni con barre verticali",
                "Analisi delle correlazioni temporali"
            ]
        )
    
    # Slide 10: Box Plots
    if os.path.exists("plots/box_plots.png"):
        add_image_slide(
            prs,
            "Box Plots: Distribuzione Features per Occupancy",
            "plots/box_plots.png",
            "Confronto delle distribuzioni con e senza occupazione"
        )
    else:
        add_content_slide(
            prs,
            "Box Plots delle Features",
            [
                "Distribuzione delle features rispetto all'occupazione",
                "Comparazione tra Occupancy=0 e Occupancy=1"
            ]
        )
    
    # Slide 11: Interpretazione Box Plots
    add_content_slide(
        prs,
        "Analisi dei Box Plots",
        [
            "Humidity: i notch si sovrappongono",
            "  - Le distribuzioni potrebbero essere simili",
            "  - Feature probabilmente meno importante per la classificazione",
            "Temperature: differenze significative",
            "  - I notch non si sovrappongono",
            "  - Feature importante per discriminare l'occupazione",
            "Light, CO2, HumidityRatio: pattern distintivi evidenti"
        ]
    )
    
    # Slide 12: Scelta del Modello
    add_content_slide(
        prs,
        "Scelta dell'Algoritmo: Random Forest",
        [
            "Perché Random Forest?",
            "  - Ensemble method robusto",
            "  - Gestisce bene features non lineari",
            "  - Riduce overfitting rispetto a singoli alberi",
            "  - Fornisce feature importance",
            "  - Ottimo bilanciamento bias-variance",
            "Caratteristiche:",
            "  - Molteplici alberi di decisione",
            "  - Voting per la predizione finale",
            "  - Bootstrap aggregating (bagging)"
        ]
    )
    
    # Slide 13: Ottimizzazione Iper-parametri
    add_content_slide(
        prs,
        "Metodologia di Ottimizzazione",
        [
            "Grid Search con 50 random_state diversi",
            "Parametri ottimizzati:",
            "  - n_estimators: [100, 200, 500]",
            "  - min_samples_split: [2, 5, 10]",
            "  - min_samples_leaf: [1, 2, 5]",
            "  - max_samples: [0.5, 0.8, 1.0]",
            "  - max_depth: [None, 10, 30]",
            "Strategia: calcolare la moda dei parametri ottimali",
            "Obiettivo: trovare i parametri più robusti e generalizzabili"
        ]
    )
    
    # Slide 14: Risultati Grid Search
    if os.path.exists("plots/frequenze_parametri_migliori_50_grid_search.png"):
        add_image_slide(
            prs,
            "Frequenze dei Parametri Migliori",
            "plots/frequenze_parametri_migliori_50_grid_search.png",
            "Distribuzione dei parametri ottimali su 50 esecuzioni diverse"
        )
    else:
        add_content_slide(
            prs,
            "Risultati Grid Search",
            [
                "Visualizzazione della frequenza dei parametri ottimali",
                "Parametri con moda evidenziata",
                "Conferma della scelta robusta degli iper-parametri"
            ]
        )
    
    # Slide 15: Parametri Ottimali
    add_content_slide(
        prs,
        "Parametri Ottimali Selezionati",
        [
            "Dai 50 esperimenti di Grid Search:",
            "I parametri con frequenza più alta (moda) sono stati:",
            "  - n_estimators: 500",
            "  - min_samples_split: 2",
            "  - min_samples_leaf: 1",
            "  - max_samples: 0.5",
            "  - max_depth: None (nessun limite)",
            "Gli iper-parametri della moda appaiono con buon distacco",
            "Indica robustezza della configurazione ottimale"
        ]
    )
    
    # Slide 16: Training del Modello
    add_content_slide(
        prs,
        "Training del Modello",
        [
            "Configurazione finale Random Forest con parametri ottimali",
            "Training eseguito su training set (75% dei dati)",
            "Utilizzo di tutti i core disponibili (n_jobs=-1)",
            "Fitting completato con successo",
            "Modello pronto per la fase di testing e validazione"
        ]
    )
    
    # Slide 17: Metriche di Valutazione
    add_content_slide(
        prs,
        "Metriche di Valutazione Utilizzate",
        [
            "Matrice di Confusione:",
            "  - True Positives, False Positives, True Negatives, False Negatives",
            "F1-Score: media armonica di Precision e Recall",
            "Recall: capacità di identificare i positivi",
            "Precision: accuratezza delle predizioni positive",
            "ROC Curve: performance al variare della soglia",
            "AUC-ROC: area sotto la curva ROC",
            "  - Misura complessiva indipendente dalla soglia"
        ]
    )
    
    # Slide 18: Risultati Testing
    add_content_slide(
        prs,
        "Risultati sul Testing Set",
        [
            "Performance sul test set (subsampling stratificato):",
            "Eccellenti risultati su tutte le metriche",
            "Alta accuratezza nella predizione",
            "Bassi falsi positivi e falsi negativi",
            "ROC-AUC molto elevata",
            "Il modello outperforma significativamente un classificatore casuale",
            "Evidenza di un ottimo lavoro di classificazione"
        ]
    )
    
    # Slide 19: Validazione Dataset Separato
    add_content_slide(
        prs,
        "Validazione su Dataset Separato (datatest2.txt)",
        [
            "Test finale su dataset completamente indipendente",
            "Simula uno scenario real-world",
            "Risultati coerenti con il testing precedente",
            "Tutti gli indici di performance confermati",
            "Dimostra la capacità di generalizzazione del modello",
            "Il modello non è overfittato sui dati di training",
            "Conferma della robustezza della soluzione"
        ]
    )
    
    # Slide 20: Feature Importance
    if os.path.exists("plots/feature_importances_random_forest.png"):
        add_image_slide(
            prs,
            "Importanza delle Features",
            "plots/feature_importances_random_forest.png",
            "Importanza relativa di ogni feature nel modello Random Forest"
        )
    else:
        add_content_slide(
            prs,
            "Feature Importance",
            [
                "Analisi dell'importanza delle features nella Random Forest",
                "Identifica quali features influenzano maggiormente la predizione"
            ]
        )
    
    # Slide 21: Interpretazione Feature Importance
    add_content_slide(
        prs,
        "Interpretazione Feature Importance",
        [
            "Le features sono classificate per importanza",
            "Light (luminosità) potrebbe essere tra le più importanti",
            "CO2 contribuisce significativamente alla classificazione",
            "Temperature gioca un ruolo importante",
            "HumidityRatio aggiunge informazioni utili",
            "Humidity, come previsto dall'analisi dei box plot, ha minor peso",
            "Conferma delle intuizioni dall'analisi esplorativa"
        ]
    )
    
    # Slide 22: Esempio Albero Decisionale
    if os.path.exists("plots/tree_random_forest.png"):
        add_image_slide(
            prs,
            "Esempio: Un Albero della Random Forest",
            "plots/tree_random_forest.png",
            "Visualizzazione di uno degli alberi decisionali della Random Forest"
        )
    else:
        add_content_slide(
            prs,
            "Struttura degli Alberi Decisionali",
            [
                "La Random Forest è composta da multipli alberi",
                "Ogni albero fa decisioni basate sulle features",
                "Struttura gerarchica di split decisionali"
            ]
        )
    
    # Slide 23: Vantaggi della Soluzione
    add_content_slide(
        prs,
        "Vantaggi della Soluzione Implementata",
        [
            "Accuratezza elevata nelle predizioni",
            "Robustezza dimostrata su dataset separati",
            "Metodologia rigorosa di ottimizzazione",
            "Validazione multipla con diverse metriche",
            "Interpretabilità attraverso feature importance",
            "Scalabilità per applicazioni real-time",
            "Facilità di deployment in sistemi esistenti"
        ]
    )
    
    # Slide 24: Possibili Applicazioni
    add_content_slide(
        prs,
        "Applicazioni Pratiche",
        [
            "Smart Building Management:",
            "  - Regolazione automatica HVAC",
            "  - Ottimizzazione consumo energetico",
            "Sicurezza:",
            "  - Monitoraggio accessi non autorizzati",
            "  - Alert in caso di anomalie",
            "Facility Management:",
            "  - Pianificazione pulizie",
            "  - Gestione spazi condivisi",
            "Analisi dati:",
            "  - Pattern di utilizzo degli spazi",
            "  - Ottimizzazione layout uffici"
        ]
    )
    
    # Slide 25: Limitazioni e Sviluppi Futuri
    add_content_slide(
        prs,
        "Limitazioni e Possibili Miglioramenti",
        [
            "Limitazioni attuali:",
            "  - Dipendenza dalla qualità dei sensori",
            "  - Necessità di calibrazione periodica",
            "  - Limitato a condizioni ambientali simili al training",
            "Sviluppi futuri:",
            "  - Integrazione di sensori aggiuntivi (es. PIR, audio)",
            "  - Transfer learning per nuovi ambienti",
            "  - Predizione del numero di occupanti",
            "  - Implementazione in edge computing devices",
            "  - Dashboard real-time per monitoraggio"
        ]
    )
    
    # Slide 26: Tecnologie Utilizzate
    add_content_slide(
        prs,
        "Stack Tecnologico",
        [
            "Linguaggio: Python 3",
            "Machine Learning: scikit-learn",
            "  - RandomForestClassifier",
            "  - GridSearchCV per ottimizzazione",
            "Data Analysis: pandas",
            "Visualizzazione:",
            "  - matplotlib per grafici",
            "  - seaborn per visualizzazioni statistiche",
            "Environment: Jupyter Notebook",
            "Version Control: Git"
        ]
    )
    
    # Slide 27: Conclusioni
    add_content_slide(
        prs,
        "Conclusioni",
        [
            "Obiettivo raggiunto con successo",
            "Modello Random Forest con ottime performance",
            "Ottimizzazione iper-parametri rigorosa e robusta",
            "Validazione multipla conferma la qualità del modello",
            "Feature importance allineata con l'analisi esplorativa",
            "Soluzione pronta per deployment in scenari reali",
            "Metodologia replicabile e scalabile",
            "Buon bilanciamento tra accuratezza e interpretabilità"
        ]
    )
    
    # Slide 28: Ringraziamenti e Q&A
    add_title_slide(
        prs,
        "Grazie per l'Attenzione!",
        "Domande?"
    )
    
    # Salva la presentazione
    output_file = "Presentazione_Occupazione_Stanza.pptx"
    prs.save(output_file)
    print(f"Presentazione creata con successo: {output_file}")
    print(f"Numero totale di slide: {len(prs.slides)}")
    
    return output_file

if __name__ == "__main__":
    output_file = create_presentation()
    print(f"\nFile creato: {output_file}")
    print("La presentazione è pronta per essere utilizzata!")
